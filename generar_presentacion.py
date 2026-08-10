import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def crear_presentacion():
    prs = Presentation()
    
    # Configurar diapositivas panorámicas (16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Colores corporativos (adaptados a RuralConecta AI)
    c_rojo = RGBColor(46, 125, 50)       # Verde rural
    c_verde = RGBColor(230, 161, 92)     # Oro tierra
    c_oscuro = RGBColor(31, 41, 55)
    c_gris_claro = RGBColor(241, 245, 249)
    c_blanco = RGBColor(255, 255, 255)
    c_gris_texto = RGBColor(75, 85, 99)
    
    # Helper: añadir fondo sólido a diapositiva
    def aplicar_fondo(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper: añadir título estándar
    def agregar_titulo(slide, texto, color=c_rojo):
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = texto
        p.font.name = 'Arial'
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = color
        
    # Helper: agregar texto de pie de página
    def agregar_footer(slide):
        footer_box = slide.shapes.add_textbox(Inches(0.75), Inches(7.0), Inches(11.83), Inches(0.3))
        tf = footer_box.text_frame
        p = tf.paragraphs[0]
        p.text = 'RuralConecta AI 🌱 Provincia de La Rioja | "Conectando Parajes, Cultivando Comunidad"'
        p.font.name = 'Arial'
        p.font.size = Pt(10)
        p.font.italic = True
        p.font.color.rgb = c_gris_texto

    # ==========================================
    # DIAPOSITIVA 1: Portada (Fondo Rojo Corporativo)
    # ==========================================
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide1 = prs.slides.add_slide(slide_layout)
    aplicar_fondo(slide1, c_rojo)
    
    # Contenedor central de título
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "🌱 Ecosistema RuralConecta AI"
    p.font.name = 'Arial'
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = c_blanco
    p.alignment = PP_ALIGN.LEFT
    
    p2 = tf.add_paragraph()
    p2.text = "Manual de Usuario: Módulo de Auditoría y SLAs"
    p2.font.name = 'Arial'
    p2.font.size = Pt(28)
    p2.font.bold = False
    p2.font.color.rgb = c_blanco
    p2.space_before = Pt(15)
    p2.alignment = PP_ALIGN.LEFT
    
    p3 = tf.add_paragraph()
    p3.text = "Herramienta analítica de monitoreo, SLAs y control de gestión rural\nProvincia de La Rioja - \"Conectando Parajes, Cultivando Comunidad\""
    p3.font.name = 'Arial'
    p3.font.size = Pt(16)
    p3.font.italic = True
    p3.font.color.rgb = c_gris_claro
    p3.space_before = Pt(40)
    p3.alignment = PP_ALIGN.LEFT

    # ==========================================
    # DIAPOSITIVA 2: Introducción al Módulo (Fondo Claro)
    # ==========================================
    slide2 = prs.slides.add_slide(slide_layout)
    aplicar_fondo(slide2, c_gris_claro)
    agregar_titulo(slide2, "📌 1. Introducción al Módulo de Auditoría")
    
    # Cuadro principal de texto
    content_box = slide2.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "El módulo de auditoría es la herramienta analítica de control interno. Registra automáticamente cada cambio de estado en la vida de una solicitud rural, permitiendo:"
    p.font.name = 'Arial'
    p.font.size = Pt(18)
    p.font.color.rgb = c_oscuro
    p.space_after = Pt(20)
    
    items = [
        ("📈 Monitoreo Operativo:", "Supervisar tiempos de respuesta y resoluciones en tiempo real."),
        ("⏳ Cuellos de Botella:", "Identificar en qué estados se acumulan retrasos (mesa de entrada vs coordinadores)."),
        ("🔍 Auditoría Forense:", "Investigar la historia completa de transiciones paso a paso por ID de caso."),
        ("👤 Desempeño del Personal:", "Evaluar la productividad y carga laboral de cada gestor de servicio."),
        ("🚨 Prevención de Demoras:", "Monitorear alertas de vencimiento para priorizar incidentes rurales críticos.")
    ]
    
    for title, desc in items:
        p_item = tf.add_paragraph()
        p_item.text = f"• {title} {desc}"
        p_item.font.name = 'Arial'
        p_item.font.size = Pt(16)
        p_item.font.color.rgb = c_oscuro
        p_item.space_before = Pt(8)
        p_item.level = 0
        
    agregar_footer(slide2)

    # ==========================================
    # DIAPOSITIVA 3: Indicadores Clave - KPIs (Fondo Claro)
    # ==========================================
    slide3 = prs.slides.add_slide(slide_layout)
    aplicar_fondo(slide3, c_gris_claro)
    agregar_titulo(slide3, "📊 2. Indicadores Clave de Rendimiento (KPIs)")
    
    # Añadir texto descriptivo
    desc_box = slide3.shapes.add_textbox(Inches(0.75), Inches(1.3), Inches(11.83), Inches(0.8))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    p = tf_desc.paragraphs[0]
    p.text = "En la parte superior del panel se presentan 5 métricas clave para evaluar la salud de la gestión vecinal:"
    p.font.name = 'Arial'
    p.font.size = Pt(16)
    p.font.color.rgb = c_oscuro
    
    # Crear tarjetas para KPIs
    kpis = [
        ("Total Reclamos", "Volumen Histórico", "Conteo acumulado de todas las solicitudes en la base de datos."),
        ("Resueltos / Cerrados", "Finalizados", "Solicitudes que han alcanzado un estado final (RESUELTO o RECHAZADO)."),
        ("Promedio Resolución", "Tiempo Medio", "Diferencia promedio en horas entre la creación y el cierre definitivo."),
        ("Cumplimiento SLA", "Efectividad %", "Porcentaje de reclamos cerrados dentro del plazo límite asignado."),
        ("Alertas SLA", "Vencimientos Activos", "Solicitudes abiertas que ya han excedido el tiempo estipulado.")
    ]
    
    card_width = Inches(2.15)
    card_height = Inches(3.8)
    spacing = Inches(0.27)
    left_start = Inches(0.75)
    top_pos = Inches(2.2)
    
    for idx, (title, subtitle, desc) in enumerate(kpis):
        # Fondo de la tarjeta (usando forma rectangular)
        left = left_start + idx * (card_width + spacing)
        shape = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top_pos, card_width, card_height)
        shape.fill.solid()
        # Destacar la de Alertas SLA en rojo
        if idx == 4:
            shape.fill.fore_color.rgb = RGBColor(254, 242, 242)
            shape.line.color.rgb = c_rojo
        else:
            shape.fill.fore_color.rgb = c_blanco
            shape.line.color.rgb = RGBColor(226, 232, 240)
            
        # Añadir texto en la tarjeta
        tf_card = shape.text_frame
        tf_card.word_wrap = True
        tf_card.margin_left = Inches(0.15)
        tf_card.margin_right = Inches(0.15)
        tf_card.margin_top = Inches(0.2)
        
        # Título
        p_title = tf_card.paragraphs[0]
        p_title.text = title
        p_title.font.name = 'Arial'
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.font.color.rgb = c_rojo if idx == 4 else c_oscuro
        p_title.alignment = PP_ALIGN.CENTER
        
        # Subtítulo
        p_sub = tf_card.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.name = 'Arial'
        p_sub.font.size = Pt(14)
        p_sub.font.bold = True
        p_sub.font.color.rgb = c_rojo if idx == 4 else c_verde
        p_sub.space_before = Pt(10)
        p_sub.alignment = PP_ALIGN.CENTER
        
        # Descripción
        p_desc = tf_card.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = 'Arial'
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = c_gris_texto
        p_desc.space_before = Pt(15)
        p_desc.alignment = PP_ALIGN.CENTER
        
    agregar_footer(slide3)

    # ==========================================
    # DIAPOSITIVA 4: Fórmulas y SLAs (Fondo Claro)
    # ==========================================
    slide4 = prs.slides.add_slide(slide_layout)
    aplicar_fondo(slide4, c_gris_claro)
    agregar_titulo(slide4, "📈 3. Fórmulas de Tiempos de Servicio (SLA)")
    
    # Explicación de SLA
    explanation_box = slide4.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(1.2))
    tf_exp = explanation_box.text_frame
    tf_exp.word_wrap = True
    p = tf_exp.paragraphs[0]
    p.text = "El cálculo del límite de tiempo se basa en la tabla de categorías (ej. Bacheo = 72 horas, Alumbrado = 24 horas, Arbolado = 48 horas). Las métricas de auditoría aplican las siguientes fórmulas en tiempo real:"
    p.font.name = 'Arial'
    p.font.size = Pt(16)
    p.font.color.rgb = c_oscuro
    
    # Caja 1: Tiempo de Resolución
    box1 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(2.9), Inches(5.6), Inches(3.6))
    box1.fill.solid()
    box1.fill.fore_color.rgb = c_blanco
    box1.line.color.rgb = RGBColor(226, 232, 240)
    tf1 = box1.text_frame
    tf1.word_wrap = True
    tf1.margin_left = Inches(0.3)
    tf1.margin_top = Inches(0.3)
    
    p_t1 = tf1.paragraphs[0]
    p_t1.text = "Tiempo de Resolución Real (Tr)"
    p_t1.font.name = 'Arial'
    p_t1.font.size = Pt(20)
    p_t1.font.bold = True
    p_t1.font.color.rgb = c_rojo
    
    p_f1 = tf1.add_paragraph()
    p_f1.text = "\nTr = FechaResolución - FechaCreación\n"
    p_f1.font.name = 'Courier New'
    p_f1.font.size = Pt(16)
    p_f1.font.bold = True
    p_f1.font.color.rgb = c_oscuro
    p_f1.alignment = PP_ALIGN.CENTER
    
    p_d1 = tf1.add_paragraph()
    p_d1.text = "Mide la diferencia de horas netas transcurridas desde el ingreso de la solicitud hasta el cierre final (RESUELTO o RECHAZADO)."
    p_d1.font.name = 'Arial'
    p_d1.font.size = Pt(13)
    p_d1.font.color.rgb = c_gris_texto
    p_d1.space_before = Pt(10)
    
    # Caja 2: Desviación de SLA
    box2 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.98), Inches(2.9), Inches(5.6), Inches(3.6))
    box2.fill.solid()
    box2.fill.fore_color.rgb = c_blanco
    box2.line.color.rgb = RGBColor(226, 232, 240)
    tf2 = box2.text_frame
    tf2.word_wrap = True
    tf2.margin_left = Inches(0.3)
    tf2.margin_top = Inches(0.3)
    
    p_t2 = tf2.paragraphs[0]
    p_t2.text = "Desviación o Retraso de SLA (Dsla)"
    p_t2.font.name = 'Arial'
    p_t2.font.size = Pt(20)
    p_t2.font.bold = True
    p_t2.font.color.rgb = c_verde
    
    p_f2 = tf2.add_paragraph()
    p_f2.text = "\nDsla = (FechaActual - FechaCreación) - SLAHoras\n"
    p_f2.font.name = 'Courier New'
    p_f2.font.size = Pt(16)
    p_f2.font.bold = True
    p_f2.font.color.rgb = c_oscuro
    p_f2.alignment = PP_ALIGN.CENTER
    
    p_d2 = tf2.add_paragraph()
    p_d2.text = "Aplica a los reclamos abiertos. Si Dsla > 0, significa que la solicitud superó el plazo máximo y activa automáticamente una Alerta SLA."
    p_d2.font.name = 'Arial'
    p_d2.font.size = Pt(13)
    p_d2.font.color.rgb = c_gris_texto
    p_d2.space_before = Pt(10)
    
    agregar_footer(slide4)

    # ==========================================
    # DIAPOSITIVA 5: Reportes Analíticos - Parte 1 (Fondo Claro)
    # ==========================================
    slide5 = prs.slides.add_slide(slide_layout)
    aplicar_fondo(slide5, c_gris_claro)
    agregar_titulo(slide5, "🛠️ 4. Reportes de Auditoría Disponibles (1/2)")
    
    reports1 = [
        ("📊 Reporte 1: Desempeño por Categoría", [
            "Muestra estadísticas agrupadas por tipo de reclamo vecinal.",
            "Visualiza: SLA establecido, conteo total, resueltos, promedio real de resolución, cumplimiento % y cantidad de casos vencidos activos."
        ]),
        ("⏳ Reporte 2: Análisis de Cuellos de Botella (Dwell Times)", [
            "Mide el tiempo de espera neto promedio de los expedientes en cada estado.",
            "Permite el diagnóstico: si hay demora en PENDIENTE el retraso es administrativo; si es en EN PROCESO, es de cuadrillas de calle."
        ]),
        ("🔍 Reporte 3: Trazabilidad Cronológica de un Reclamo", [
            "Auditoría forense individualizada al ingresar el identificador único (ID).",
            "Muestra metadatos generales y una línea de tiempo paso a paso con fecha, hora, operador responsable y rol del cambio."
        ])
    ]
    
    for idx, (title, bullets) in enumerate(reports1):
        # Crear contenedor para cada reporte
        top = Inches(1.5 + idx * 1.8)
        shape = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), top, Inches(11.83), Inches(1.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = c_blanco
        shape.line.color.rgb = RGBColor(226, 232, 240)
        
        tf_rep = shape.text_frame
        tf_rep.word_wrap = True
        tf_rep.margin_left = Inches(0.25)
        tf_rep.margin_top = Inches(0.15)
        
        p_t = tf_rep.paragraphs[0]
        p_t.text = title
        p_t.font.name = 'Arial'
        p_t.font.size = Pt(18)
        p_t.font.bold = True
        p_t.font.color.rgb = c_rojo
        
        for b in bullets:
            p_b = tf_rep.add_paragraph()
            p_b.text = f"• {b}"
            p_b.font.name = 'Arial'
            p_b.font.size = Pt(13)
            p_b.font.color.rgb = c_oscuro
            p_b.space_before = Pt(4)
            
    agregar_footer(slide5)

    # ==========================================
    # DIAPOSITIVA 6: Reportes Analíticos - Parte 2 (Fondo Claro)
    # ==========================================
    slide6 = prs.slides.add_slide(slide_layout)
    aplicar_fondo(slide6, c_gris_claro)
    agregar_titulo(slide6, "🛠️ 4. Reportes de Auditoría Disponibles (2/2)")
    
    reports2 = [
        ("⚠️ Reporte 4: Alertas de Desviación y SLA Excedido", [
            "Lista priorizada de reclamos abiertos fuera de tiempo, ordenada por prioridad y mayor retraso acumulado.",
            "Facilita la reasignación inmediata a gestores y el seguimiento del plazo de entrega."
        ]),
        ("👤 Reporte 5: Desempeño de Gestores", [
            "Mide la productividad, carga laboral y eficiencia de cada operador municipal asignado.",
            "Métricas: Casos activos, cerrados exitosamente, promedio de resolución y % de casos cumplidos a tiempo."
        ]),
        ("⚙️ Reportes 6 y 7: Testing y Visualización Gráfica", [
            "Opción 6: Permite simular transiciones de prueba coherentes para validar el sistema analítico (borra historial previo).",
            "Opción 7: Genera gráficos interactivos sobre reclamos por barrio, categoría, estados y prioridades en Streamlit."
        ])
    ]
    
    for idx, (title, bullets) in enumerate(reports2):
        top = Inches(1.5 + idx * 1.8)
        shape = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), top, Inches(11.83), Inches(1.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = c_blanco
        shape.line.color.rgb = RGBColor(226, 232, 240)
        
        tf_rep = shape.text_frame
        tf_rep.word_wrap = True
        tf_rep.margin_left = Inches(0.25)
        tf_rep.margin_top = Inches(0.15)
        
        p_t = tf_rep.paragraphs[0]
        p_t.text = title
        p_t.font.name = 'Arial'
        p_t.font.size = Pt(18)
        p_t.font.bold = True
        p_t.font.color.rgb = c_rojo
        
        for b in bullets:
            p_b = tf_rep.add_paragraph()
            p_b.text = f"• {b}"
            p_b.font.name = 'Arial'
            p_b.font.size = Pt(13)
            p_b.font.color.rgb = c_oscuro
            p_b.space_before = Pt(4)
            
    agregar_footer(slide6)

    # ==========================================
    # DIAPOSITIVA 7: Ciclo de Vida y Roles (Fondo Claro)
    # ==========================================
    slide7 = prs.slides.add_slide(slide_layout)
    aplicar_fondo(slide7, c_gris_claro)
    agregar_titulo(slide7, "🔄 5. Ciclo de Vida del Reclamo y Roles")
    
    # Título sección 1
    t1_box = slide7.shapes.add_textbox(Inches(0.75), Inches(1.3), Inches(5.6), Inches(0.5))
    t1_box.text_frame.paragraphs[0].text = "Flujo Operativo de Estados:"
    t1_box.text_frame.paragraphs[0].font.name = 'Arial'
    t1_box.text_frame.paragraphs[0].font.size = Pt(18)
    t1_box.text_frame.paragraphs[0].font.bold = True
    t1_box.text_frame.paragraphs[0].font.color.rgb = c_oscuro
    
    # Caja para flujo de estados
    box_flow = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.9), Inches(5.6), Inches(4.8))
    box_flow.fill.solid()
    box_flow.fill.fore_color.rgb = c_blanco
    box_flow.line.color.rgb = RGBColor(226, 232, 240)
    tf_flow = box_flow.text_frame
    tf_flow.word_wrap = True
    tf_flow.margin_left = Inches(0.3)
    tf_flow.margin_top = Inches(0.3)
    
    flow_steps = [
        ("1. PENDIENTE", "Ingresado por el vecino o analista a la base de datos."),
        ("2. EN REVISION", "El gestor valida datos, asigna nivel de prioridad y operador."),
        ("3. EN PROCESO", "Caso derivado a cuadrillas de trabajo urbano en la calle."),
        ("4. RESUELTO", "Cierre definitivo por solución técnica del inconveniente."),
        ("5. RECHAZADO", "Cierre final de caso duplicado, inválido o fuera de rango.")
    ]
    
    for idx, (step, desc) in enumerate(flow_steps):
        p_step = tf_flow.add_paragraph() if idx > 0 else tf_flow.paragraphs[0]
        p_step.text = step
        p_step.font.name = 'Arial'
        p_step.font.size = Pt(15)
        p_step.font.bold = True
        p_step.font.color.rgb = c_rojo
        if idx > 0:
            p_step.space_before = Pt(12)
            
        p_desc = tf_flow.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = 'Arial'
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = c_gris_texto
        p_desc.space_before = Pt(2)

    # Título sección 2
    t2_box = slide7.shapes.add_textbox(Inches(6.98), Inches(1.3), Inches(5.6), Inches(0.5))
    t2_box.text_frame.paragraphs[0].text = "Roles y Responsabilidades:"
    t2_box.text_frame.paragraphs[0].font.name = 'Arial'
    t2_box.text_frame.paragraphs[0].font.size = Pt(18)
    t2_box.text_frame.paragraphs[0].font.bold = True
    t2_box.text_frame.paragraphs[0].font.color.rgb = c_oscuro
    
    # Caja para roles
    box_roles = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.98), Inches(1.9), Inches(5.6), Inches(4.8))
    box_roles.fill.solid()
    box_roles.fill.fore_color.rgb = c_blanco
    box_roles.line.color.rgb = RGBColor(226, 232, 240)
    tf_roles = box_roles.text_frame
    tf_roles.word_wrap = True
    tf_roles.margin_left = Inches(0.3)
    tf_roles.margin_top = Inches(0.3)
    
    roles_info = [
        ("👤 Vecino (ID 1):", "Ingresa el reclamo y lo sitúa en PENDIENTE. Su comentario es inmutable para asegurar transparencia."),
        ("💼 Gestor (ID 2):", "Modifica el estado del caso operativamente. El historial registra sus acciones con nombre y apellido."),
        ("📈 Analista (ID 3):", "Supervisa las alertas de desviación global, reasigna casos y descarga reportes analíticos."),
        ("⚙️ Sistema (Automático):", "Genera logs en transiciones de testing o ante fallas en los disparadores relacionales.")
    ]
    
    for idx, (role, desc) in enumerate(roles_info):
        p_role = tf_roles.add_paragraph() if idx > 0 else tf_roles.paragraphs[0]
        p_role.text = role
        p_role.font.name = 'Arial'
        p_role.font.size = Pt(15)
        p_role.font.bold = True
        p_role.font.color.rgb = c_verde
        if idx > 0:
            p_role.space_before = Pt(12)
            
        p_desc = tf_roles.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = 'Arial'
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = c_gris_texto
        p_desc.space_before = Pt(2)
        
    agregar_footer(slide7)

    # ==========================================
    # DIAPOSITIVA 8: Buenas Prácticas (Fondo Claro)
    # ==========================================
    slide8 = prs.slides.add_slide(slide_layout)
    aplicar_fondo(slide8, c_gris_claro)
    agregar_titulo(slide8, "💡 6. Buenas Prácticas de Gestión de Auditoría")
    
    # Contenido principal
    content_box = slide8.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Para maximizar la efectividad del sistema, los coordinadores del medio rural deben seguir este plan de buenas prácticas:"
    p.font.name = 'Arial'
    p.font.size = Pt(18)
    p.font.color.rgb = c_oscuro
    p.space_after = Pt(25)
    
    practices = [
        ("📅 Monitoreo Diario de Alertas SLA:", "Iniciar la jornada revisando el Reporte 4 (Alertas de Desviación) para reasignar tareas críticas con mayor retraso acumulado."),
        ("⏳ Control Semanal de Cuellos de Botella:", "Supervisar el Reporte 2 para identificar incrementos en la permanencia en PENDIENTE o EN PROCESO (falta de respuesta en parajes)."),
        ("🔍 Auditoría Forense por ID:", "Utilizar el Reporte 3 ante consultas de ciudadanos o coordinadores locales para documentar con fecha, hora y responsable cada etapa de la gestión."),
        ("⚙️ Simulación Segura:", "Utilizar la herramienta de simulación (Reporte 6) únicamente en entornos de prueba, ya que limpia el historial real de auditoría.")
    ]
    
    for idx, (title, desc) in enumerate(practices):
        p_t = tf.add_paragraph()
        p_t.text = title
        p_t.font.name = 'Arial'
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = c_rojo
        p_t.space_before = Pt(15)
        
        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.name = 'Arial'
        p_d.font.size = Pt(14)
        p_d.font.color.rgb = c_oscuro
        p_d.space_before = Pt(3)
        p_d.margin_left = Inches(0.3)
        
    agregar_footer(slide8)

    # Guardar presentación
    output_filename = "Presentacion_Auditoria_RuralConecta.pptx"
    prs.save(output_filename)
    print(f"Presentation saved successfully as {output_filename}!")

if __name__ == '__main__':
    crear_presentacion()
