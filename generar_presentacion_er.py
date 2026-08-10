import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def crear_presentacion_er():
    prs = Presentation()
    
    # Configurar diapositivas panorámicas (16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Colores corporativos (adaptados a RuralConecta AI)
    c_rojo = RGBColor(46, 125, 50)       # Verde rural
    c_verde = RGBColor(230, 161, 92)     # Oro tierra
    c_oscuro = RGBColor(31, 41, 55)
    c_gris_claro = RGBColor(241, 245, 249)
    c_blanco = RGBColor(255, 255, 255)
    c_gris_texto = RGBColor(75, 85, 99)
    
    # Helper: añadir fondo sólido a diapositiva
    def aplicar_fondo(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper: añadir título estándar
    def agregar_titulo(slide, texto, color=c_rojo):
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = texto
        p.font.name = 'Arial'
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = color
        
    # Helper: agregar texto de pie de página
    def agregar_footer(slide):
        footer_box = slide.shapes.add_textbox(Inches(0.75), Inches(7.0), Inches(11.83), Inches(0.3))
        tf = footer_box.text_frame
        p = tf.paragraphs[0]
        p.text = 'Modelo de Entidad-Relación | RuralConecta AI 🌱 Provincia de La Rioja'
        p.font.name = 'Arial'
        p.font.size = Pt(10)
        p.font.italic = True
        p.font.color.rgb = c_gris_texto

    # ==========================================
    # DIAPOSITIVA 1: Portada (Fondo Verde Municipal)
    # ==========================================
    slide_layout = prs.slide_layouts[6] # Blank
    slide1 = prs.slides.add_slide(slide_layout)
    aplicar_fondo(slide1, c_verde)
    
    # Contenedor central de título
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "🌱 Modelo de Entidad-Relación (MER)"
    p.font.name = 'Arial'
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = c_blanco
    p.alignment = PP_ALIGN.LEFT
    
    p2 = tf.add_paragraph()
    p2.text = "Arquitectura de Base de Datos y Trazabilidad Relacional"
    p2.font.name = 'Arial'
    p2.font.size = Pt(28)
    p2.font.bold = False
    p2.font.color.rgb = c_blanco
    p2.space_before = Pt(15)
    p2.alignment = PP_ALIGN.LEFT
    
    p3 = tf.add_paragraph()
    p3.text = "Basado en el esquema oficial de RuralConecta AI (ruralconecta.db)\nProvincia de La Rioja - \"Conectando Parajes, Cultivando Comunidad\""
    p3.font.name = 'Arial'
    p3.font.size = Pt(16)
    p3.font.italic = True
    p3.font.color.rgb = c_gris_claro
    p3.space_before = Pt(40)
    p3.alignment = PP_ALIGN.LEFT

    # ==========================================
    # DIAPOSITIVA 2: Vista General (Fondo Claro)
    # ==========================================
    slide2 = prs.slides.add_slide(slide_layout)
    aplicar_fondo(slide2, c_gris_claro)
    agregar_titulo(slide2, "📌 1. Vista General de la Base de Datos")
    
    content_box = slide2.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "La base de datos (`ruralconecta.db`) está diseñada bajo el motor relacional SQLite con integridad referencial activa (PRAGMA foreign_keys = ON). Se compone de 16 tablas divididas en 4 áreas operativas:"
    p.font.name = 'Arial'
    p.font.size = Pt(18)
    p.font.color.rgb = c_oscuro
    p.space_after = Pt(20)
    
    areas = [
        ("🔐 1. Seguridad y Accesos:", "Control de permisos de roles (Ciudadano Rural vs Gestor) y credenciales seguras (SHA-256)."),
        ("📥 2. Flujo Operativo de Solicitudes:", "Estructura jerárquica de categorías, subcategorías, parajes y solicitudes."),
        ("🌱 3. Auditoría y SLAs:", "Bitácora inmutable de transiciones de estados que registra fecha, hora y responsable."),
        ("🤖 4. Módulos Analíticos e Inteligencia Artificial:", "Almacenamiento de predicciones NLP, análisis de anomalías, encuestas y registro meteorológico.")
    ]
    
    for title, desc in areas:
        p_item = tf.add_paragraph()
        p_item.text = f"• {title} {desc}"
        p_item.font.name = 'Arial'
        p_item.font.size = Pt(16)
        p_item.font.color.rgb = c_oscuro
        p_item.space_before = Pt(10)
        
    agregar_footer(slide2)

    # ==========================================
    # DIAPOSITIVA 3: Imagen del Modelo de Base de Datos (Fondo Oscuro)
    # ==========================================
    slide3 = prs.slides.add_slide(slide_layout)
    aplicar_fondo(slide3, c_oscuro)
    agregar_titulo(slide3, "📊 2. Diagrama de Entidad-Relación Visual", c_blanco)
    
    # Ruta de la imagen copiada en el workspace
    image_path = "modelo_base_datos.png"
    if os.path.exists(image_path):
        # Insertar imagen centrada
        # Ancho slide: 13.33, Alto: 7.5. Título ocupa hasta Y=1.3.
        slide3.shapes.add_picture(image_path, Inches(1.91), Inches(1.5), Inches(9.5), Inches(5.0))
    else:
        # Texto de fallback si no encuentra la imagen
        err_box = slide3.shapes.add_textbox(Inches(2.0), Inches(3.0), Inches(9.33), Inches(2.0))
        p = err_box.text_frame.paragraphs[0]
        p.text = "[⚠️ Imagen de modelo_base_datos.png no encontrada en el directorio]"
        p.font.name = 'Arial'
        p.font.size = Pt(20)
        p.font.color.rgb = c_rojo
        p.alignment = PP_ALIGN.CENTER
        
    agregar_footer(slide3)

    # ==========================================
    # DIAPOSITIVA 4: Seguridad y Accesos (Fondo Claro)
    # ==========================================
    slide4 = prs.slides.add_slide(slide_layout)
    aplicar_fondo(slide4, c_gris_claro)
    agregar_titulo(slide4, "🔐 3. Módulo de Seguridad y Usuarios")
    
    # Añadir texto explicativo
    desc_box = slide4.shapes.add_textbox(Inches(0.75), Inches(1.3), Inches(11.83), Inches(0.8))
    tf_desc = desc_box.text_frame
    p = tf_desc.paragraphs[0]
    p.text = "Estructura que maneja la autenticación y el control de accesos diferencial en la app:"
    p.font.name = 'Arial'
    p.font.size = Pt(16)
    p.font.color.rgb = c_oscuro
    
    tables_sec = [
        ("roles", "Define los perfiles del sistema.", "• id (PK)\n• nombre (Vecino, De Gestión)"),
        ("permisos", "Catálogo de permisos de seguridad.", "• id (PK)\n• nombre (CREAR_RECLAMO, GESTIONAR_USUARIOS, etc.)\n• descripcion"),
        ("roles_permisos", "Tabla intermedia de muchos-a-muchos.", "• rol_id (PK, FK)\n• permiso_id (PK, FK)\n• Relación CASCADE al borrar roles/permisos."),
        ("usuarios", "Registro de cuentas de acceso.", "• id (PK)\n• nombre, apellido, dni (único), cuil (único)\n• clave (SHA-256)\n• rol_id (FK)")
    ]
    
    card_width = Inches(2.7)
    card_height = Inches(4.0)
    spacing = Inches(0.3)
    left_start = Inches(0.75)
    top_pos = Inches(2.2)
    
    for idx, (name, purpose, fields) in enumerate(tables_sec):
        left = left_start + idx * (card_width + spacing)
        shape = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top_pos, card_width, card_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = c_blanco
        shape.line.color.rgb = RGBColor(226, 232, 240)
        
        tf_card = shape.text_frame
        tf_card.word_wrap = True
        tf_card.margin_left = Inches(0.15)
        tf_card.margin_top = Inches(0.2)
        
        p_name = tf_card.paragraphs[0]
        p_name.text = f"📋 {name}"
        p_name.font.name = 'Arial'
        p_name.font.size = Pt(18)
        p_name.font.bold = True
        p_name.font.color.rgb = c_rojo
        
        p_purp = tf_card.add_paragraph()
        p_purp.text = purpose
        p_purp.font.name = 'Arial'
        p_purp.font.size = Pt(12)
        p_purp.font.italic = True
        p_purp.font.color.rgb = c_gris_texto
        p_purp.space_before = Pt(8)
        
        p_fields = tf_card.add_paragraph()
        p_fields.text = f"\nAtributos:\n{fields}"
        p_fields.font.name = 'Arial'
        p_fields.font.size = Pt(11)
        p_fields.font.color.rgb = c_oscuro
        p_fields.space_before = Pt(10)
        
    agregar_footer(slide4)

    # ==========================================
    # DIAPOSITIVA 5: Estructura de Reclamos (Fondo Claro)
    # ==========================================
    slide5 = prs.slides.add_slide(slide_layout)
    aplicar_fondo(slide5, c_gris_claro)
    agregar_titulo(slide5, "📥 4. Estructura y Registro de Reclamos")
    
    # Cuadro de flujo operativo
    box_sol = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.5), Inches(5.6), Inches(5.0))
    box_sol.fill.solid()
    box_sol.fill.fore_color.rgb = c_blanco
    box_sol.line.color.rgb = RGBColor(226, 232, 240)
    tf_sol = box_sol.text_frame
    tf_sol.word_wrap = True
    tf_sol.margin_left = Inches(0.3)
    tf_sol.margin_top = Inches(0.3)
    
    p_ts = tf_sol.paragraphs[0]
    p_ts.text = "Tabla solicitudes (Núcleo Relacional)"
    p_ts.font.name = 'Arial'
    p_ts.font.size = Pt(20)
    p_ts.font.bold = True
    p_ts.font.color.rgb = c_rojo
    
    p_ds = tf_sol.add_paragraph()
    p_ds.text = "\nCentraliza todos los incidentes públicos. Sus atributos clave son:\n" \
                "• id: Número de Gestión Único.\n" \
                "• comentario: Problemática del vecino.\n" \
                "• prioridad: Urgencia operativa (ALTA, MEDIA, BAJA).\n" \
                "• score_sentimiento y urgencia_nlp: Métricas obtenidas mediante inteligencia artificial (NLP).\n" \
                "• fecha_creacion y fecha_resolucion: Marcas de tiempo de control de SLA."
    p_ds.font.name = 'Arial'
    p_ds.font.size = Pt(14)
    p_ds.font.color.rgb = c_oscuro
    p_ds.space_before = Pt(5)

    # Tablas de apoyo
    box_apoyo = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.98), Inches(1.5), Inches(5.6), Inches(5.0))
    box_apoyo.fill.solid()
    box_apoyo.fill.fore_color.rgb = c_blanco
    box_apoyo.line.color.rgb = RGBColor(226, 232, 240)
    tf_ap = box_apoyo.text_frame
    tf_ap.word_wrap = True
    tf_ap.margin_left = Inches(0.3)
    tf_ap.margin_top = Inches(0.3)
    
    p_ta = tf_ap.paragraphs[0]
    p_ta.text = "Entidades de Clasificación y Apoyo"
    p_ta.font.name = 'Arial'
    p_ta.font.size = Pt(20)
    p_ta.font.bold = True
    p_ta.font.color.rgb = c_verde
    
    p_da = tf_ap.add_paragraph()
    p_da.text = "\n• categorias: Tipos de reclamos. Contiene sla_horas (tiempo límite) y estacionalidad_alta.\n\n" \
                "• subcategorias: Desglose detallado del incidente. Vinculada a una categoría.\n\n" \
                "• barrios: Geolocalización municipal. Divide a la ciudad en zonas (Centro, Norte, Sur, Este, Oeste).\n\n" \
                "• estados_solicitud: Estados operativos habilitados (PENDIENTE, EN REVISION, EN PROCESO, RESUELTO, RECHAZADO)."
    p_da.font.name = 'Arial'
    p_da.font.size = Pt(14)
    p_da.font.color.rgb = c_oscuro
    p_da.space_before = Pt(5)
    
    agregar_footer(slide5)

    # ==========================================
    # DIAPOSITIVA 6: Trazabilidad y SLAs (Fondo Claro)
    # ==========================================
    slide6 = prs.slides.add_slide(slide_layout)
    aplicar_fondo(slide6, c_gris_claro)
    agregar_titulo(slide6, "🏛️ 5. Trazabilidad de Auditoría e Integridad SLA")
    
    # Cuadro explicativo izquierdo
    info_box = slide6.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.6), Inches(5.0))
    tf_info = info_box.text_frame
    tf_info.word_wrap = True
    
    p = tf_info.paragraphs[0]
    p.text = "La tabla historial_estados es la bitácora inmutable del sistema analítico de auditoría. Registra cada transición operativa:"
    p.font.name = 'Arial'
    p.font.size = Pt(16)
    p.font.color.rgb = c_oscuro
    p.space_after = Pt(15)
    
    items_traz = [
        ("Historial Forense:", "Permite rastrear quién cambió de estado una solicitud, cuándo y qué rol tenía."),
        ("Cálculo de Dwell Times:", "Mide de forma exacta el tiempo de permanencia de un reclamo en cada estado (por ejemplo, retrasos en revisión vs en calle)."),
        ("Auditoría de SLAs:", "Almacena los datos clave que validan si el reclamo se resolvió a tiempo o si superó el plazo máximo estipulado por su categoría.")
    ]
    for title, desc in items_traz:
        p_i = tf_info.add_paragraph()
        p_i.text = f"• {title} {desc}"
        p_i.font.name = 'Arial'
        p_i.font.size = Pt(14)
        p_i.font.color.rgb = c_oscuro
        p_i.space_before = Pt(8)
        
    # Esquema de la tabla (derecha)
    box_hist = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.98), Inches(1.5), Inches(5.6), Inches(5.0))
    box_hist.fill.solid()
    box_hist.fill.fore_color.rgb = c_blanco
    box_hist.line.color.rgb = RGBColor(226, 232, 240)
    tf_hist = box_hist.text_frame
    tf_hist.word_wrap = True
    tf_hist.margin_left = Inches(0.3)
    tf_hist.margin_top = Inches(0.3)
    
    p_th = tf_hist.paragraphs[0]
    p_th.text = "Esquema: historial_estados"
    p_th.font.name = 'Arial'
    p_th.font.size = Pt(20)
    p_th.font.bold = True
    p_th.font.color.rgb = c_rojo
    
    p_dh = tf_hist.add_paragraph()
    p_dh.text = "\nAtributos y Claves:\n\n" \
                "• id (PK): Identificador autoincremental.\n" \
                "• solicitud_id (FK): Enlace al reclamo auditado.\n" \
                "• estado_anterior_id (FK): Estado de origen.\n" \
                "• estado_nuevo_id (FK): Estado de destino.\n" \
                "• usuario_id (FK): Operador municipal responsable del cambio.\n" \
                "• fecha_cambio (DATETIME): Marca de tiempo del cambio de estado."
    p_dh.font.name = 'Arial'
    p_dh.font.size = Pt(14)
    p_dh.font.color.rgb = c_oscuro
    p_dh.space_before = Pt(10)
    
    agregar_footer(slide6)

    # ==========================================
    # DIAPOSITIVA 7: Analítica e Inteligencia Artificial (Fondo Claro)
    # ==========================================
    slide7 = prs.slides.add_slide(slide_layout)
    aplicar_fondo(slide7, c_gris_claro)
    agregar_titulo(slide7, "🤖 6. Inteligencia Artificial y Reportes Analíticos")
    
    # Cuadro descriptivo
    desc_box = slide7.shapes.add_textbox(Inches(0.75), Inches(1.3), Inches(11.83), Inches(0.8))
    tf_desc = desc_box.text_frame
    p = tf_desc.paragraphs[0]
    p.text = "Tablas estratégicas destinadas a almacenar el procesamiento avanzado del sistema:"
    p.font.name = 'Arial'
    p.font.size = Pt(16)
    p.font.color.rgb = c_oscuro
    
    tables_an = [
        ("etiquetas_ia", "Keywords NLP extraídas del reclamo para identificar problemáticas y agruparlas.", "• id (PK)\n• solicitud_id (FK)\n• keyword\n• confianza_ia (REAL)"),
        ("encuestas_satisfaccion", "Medición de calidad del servicio. Registra el feedback del vecino.", "• id (PK)\n• solicitud_id (FK)\n• puntuacion (1 a 5)\n• comentario_vecino\n• fecha_encuesta"),
        ("alertas_anomalias", "Picos atípicos y sobrecargas de reclamos en categorías críticas.", "• id (PK)\n• fecha_deteccion\n• tipo_anomalia\n• categoria_id (FK)\n• severidad"),
        ("zonas_calientes", "Mapea focos geográficos recurrentes de problemas urbanos.", "• id (PK)\n• categoria_id (FK)\n• barrio_id (FK)\n• coordenadas y recurrencia")
    ]
    
    card_width = Inches(2.7)
    card_height = Inches(4.0)
    spacing = Inches(0.3)
    left_start = Inches(0.75)
    top_pos = Inches(2.2)
    
    for idx, (name, purpose, fields) in enumerate(tables_an):
        left = left_start + idx * (card_width + spacing)
        shape = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top_pos, card_width, card_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = c_blanco
        shape.line.color.rgb = RGBColor(226, 232, 240)
        
        tf_card = shape.text_frame
        tf_card.word_wrap = True
        tf_card.margin_left = Inches(0.15)
        tf_card.margin_top = Inches(0.2)
        
        p_name = tf_card.paragraphs[0]
        p_name.text = f"🤖 {name}"
        p_name.font.name = 'Arial'
        p_name.font.size = Pt(17)
        p_name.font.bold = True
        p_name.font.color.rgb = c_verde
        
        p_purp = tf_card.add_paragraph()
        p_purp.text = purpose
        p_purp.font.name = 'Arial'
        p_purp.font.size = Pt(12)
        p_purp.font.italic = True
        p_purp.font.color.rgb = c_gris_texto
        p_purp.space_before = Pt(8)
        
        p_fields = tf_card.add_paragraph()
        p_fields.text = f"\nAtributos:\n{fields}"
        p_fields.font.name = 'Arial'
        p_fields.font.size = Pt(11)
        p_fields.font.color.rgb = c_oscuro
        p_fields.space_before = Pt(10)
        
    agregar_footer(slide7)

    # ==========================================
    # DIAPOSITIVA 8: Buenas Prácticas e Integridad (Fondo Claro)
    # ==========================================
    slide8 = prs.slides.add_slide(slide_layout)
    aplicar_fondo(slide8, c_gris_claro)
    agregar_titulo(slide8, "💡 7. Integridad de Base de Datos y Buenas Prácticas")
    
    # Cuadro principal
    content_box = slide8.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Recomendaciones técnicas para mantener la integridad de la base de datos de RuralConecta:"
    p.font.name = 'Arial'
    p.font.size = Pt(18)
    p.font.color.rgb = c_oscuro
    p.space_after = Pt(20)
    
    practices = [
        ("🔑 Activar Claves Foráneas (SQLite):", "Siempre ejecutar `PRAGMA foreign_keys = ON;` al iniciar cualquier conexión de base de datos en Python para forzar las restricciones relacionales."),
        ("🗑️ Eliminaciones en Cascada (ON DELETE CASCADE):", "Configurada en las tablas intermedias como `roles_permisos` para asegurar que al borrar un rol, se depuren automáticamente sus privilegios y evitar registros huérfanos."),
        ("🛡️ Transparencia de Auditoría:", "No se debe permitir la modificación del campo `comentario` en `solicitudes` una vez ingresado. Cualquier edición de estado debe registrarse de forma inmutable en `historial_estados`."),
        ("📊 Contextualización Meteorológica:", "Utilizar cruces relacionales por fecha con `registro_climatico` para identificar si picos atípicos de solicitudes de Caminos Rurales coinciden con temporales de lluvia de gran volumen.")
    ]
    
    for title, desc in practices:
        p_t = tf.add_paragraph()
        p_t.text = title
        p_t.font.name = 'Arial'
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = c_rojo
        p_t.space_before = Pt(12)
        
        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.name = 'Arial'
        p_d.font.size = Pt(14)
        p_d.font.color.rgb = c_oscuro
        p_d.space_before = Pt(3)
        p_d.margin_left = Inches(0.3)
        
    agregar_footer(slide8)

    # Guardar presentación
    output_filename = "Presentacion_ER_RuralConecta.pptx"
    prs.save(output_filename)
    print(f"Presentation saved successfully as {output_filename}!")

if __name__ == '__main__':
    crear_presentacion_er()
